# from powerbook import Powerbook
import os
import pptx
import tempfile
import matplotlib
import warnings


_NOTES_SECTION_DELIM = "--- Do not edit below this line ---"


def get_indent_level(text: str, indent_spaces: int = 4):
    return int((len(text) - len(text.lstrip())) / indent_spaces)


def _markdown_into_paragraphs(markdown: str) -> list:
    markdown_lines = markdown.strip().split("\n")
    first_line = markdown_lines[0]

    results = [first_line]
    for line in markdown_lines[1:]:
        if line.strip():
            results.append(
                (
                    (
                        line.strip()[2:]
                        if (
                            line.strip().startswith("- ")
                            or line.strip().startswith("* ")
                        )
                        else line.strip()
                    ),
                    get_indent_level(line),
                )
            )
        else:
            results.append(("", 0))
    return results


class Powerbook:
    def __init__(self, powerpoint_filepath: str = None) -> None:
        self._fpath = None
        if powerpoint_filepath:
            self._fpath = os.path.expanduser(powerpoint_filepath)
            self._on_disk = True
        else:
            self._on_disk = False

        self._doc = pptx.Presentation(
            self._fpath if self._fpath and os.path.exists(self._fpath) else None
        )
        self._save_on_operation_complete = True
        self._warn_dpi = True

    @property
    def slides(self):
        return self._doc.slides

    def _parse_notes_slots(self, raw_slots_text: str) -> list:
        # returns [(internal name, slide element)]
        slots = []
        for line in raw_slots_text.split("\n"):
            if line.strip():
                name, element = line.split("\t")
                slots.append((name, element))
        return slots

    @property
    def slots(self):
        slots = []
        # Go through the list of slides and get all slides
        # where the notes section has a '--- Do not edit below this line ---'
        for slide in self.slides:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if _NOTES_SECTION_DELIM not in notes:
                    continue
                raw_slots_text = notes.split(_NOTES_SECTION_DELIM)[-1]
                slide_slots = self._parse_notes_slots(raw_slots_text)
                for name, element in slide_slots:
                    slots.append((name, element, slide))
        return slots

    def _has_changed(self):
        if self._save_on_operation_complete and self._on_disk:
            self.save(self._fpath)

    def save(self, fpath: str = None):
        fpath = fpath or self._fpath
        self._doc.save(fpath)

    def add_title_slide(self, title: str, subtitle: str):
        title_slide_layout = self._doc.slide_layouts[0]
        slide = self._doc.slides.add_slide(title_slide_layout)
        _title = slide.shapes.title
        _subtitle = slide.placeholders[1]

        _title.text = title
        _subtitle.text = subtitle

        self._has_changed()

    def add_text_slide(self, title: str, text: str):
        bullet_slide_layout = self._doc.slide_layouts[1]

        slide = self._doc.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title

        text_lines = _markdown_into_paragraphs(text)

        tf = body_shape.text_frame
        tf.text = text_lines[0]

        for line, indent_level in text_lines[1:]:
            p = tf.add_paragraph()

            # if "||" in line:
            #     is_slot = False
            #     for seg in line.split("||"):
            #         if is_slot:
            #             run = p.add_run()
            #             run.text = seg
            #         else:
            #             run = p.add_run()
            #             run.text = seg
            #         is_slot = not is_slot
            p.text = line
            p.level = indent_level

        self._has_changed()

    def add_two_content_slide(self, title: str, left: str, right: str):
        layout = self._doc.slide_layouts[3]

        slide = self._doc.slides.add_slide(layout)
        shapes = slide.shapes

        title_shape = shapes.title

        left_shape = shapes.placeholders[1]
        title_shape.text = title
        text_lines = _markdown_into_paragraphs(left)

        tf = left_shape.text_frame
        tf.text = text_lines[0]

        for line, indent_level in text_lines[1:]:
            p = tf.add_paragraph()

            p.text = line
            p.level = indent_level

        # deal with matplotlib figs:

        right_shape = shapes.placeholders[2]
        if isinstance(right, matplotlib.figure.Figure):
            if right.dpi < 100 and self._warn_dpi:
                warnings.warn(
                    Warning(
                        f"The resolution of image {right} is {right.dpi}, which may not be high enough quality for this presentation."
                    )
                )
            z = tempfile.NamedTemporaryFile(suffix=".png")
            right.savefig(z.name)
            z.seek(0)

            pic = shapes.add_picture(
                z.name,
                left=right_shape.left,
                top=right_shape.top,
                width=right_shape.width,
            )
            self._has_changed()
        else:
            if os.path.exists(right):

                pic = shapes.add_picture(
                    os.path.expanduser(right),
                    left=right_shape.left,
                    top=right_shape.top,
                    width=right_shape.width,
                )
            else:
                text_lines = _markdown_into_paragraphs(right)

                tf = right_shape.text_frame
                tf.text = text_lines[0]

                for line, indent_level in text_lines[1:]:
                    p = tf.add_paragraph()

                    p.text = line
                    p.level = indent_level

            self._has_changed()

    def add_image_slide(self, title: str, image: str):
        layout = self._doc.slide_layouts[1]

        slide = self._doc.slides.add_slide(layout)
        shapes = slide.shapes

        title_shape = shapes.title
        title_shape.text = title

        # deal with matplotlib figs:

        shape = shapes.placeholders[1]
        if isinstance(image, matplotlib.figure.Figure):
            if image.dpi < 100 and self._warn_dpi:
                warnings.warn(
                    Warning(
                        f"The resolution of image {image} is {image.dpi}, which may not be high enough quality for this presentation."
                    )
                )
            z = tempfile.NamedTemporaryFile(suffix=".png")
            image.savefig(z.name)
            z.seek(0)

            pic = shapes.add_picture(
                z.name, left=shape.left, top=shape.top, width=shape.width
            )
        elif os.path.exists(image):

            pic = shapes.add_picture(
                os.path.expanduser(image),
                left=shape.left,
                top=shape.top,
                width=shape.width,
            )
        else:
            raise ValueError(f"Can't process image {image}")

        self._has_changed()